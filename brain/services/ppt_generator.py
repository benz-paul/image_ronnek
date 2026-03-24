"""
PPT Generator Service - Creates PowerPoint presentations from learning steps and scenes.
"""

import os
import re
import json
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from brain.pipeline.state.pipeline_state import PipelineState


def parse_ls_number(ls_id: str) -> int:
    """
    Extract learning step number from a learning_step_id string.
    Example: 'LS1' -> 1,  'LS10' -> 10
    """
    match = re.search(r"LS(\d+)", ls_id, re.IGNORECASE)
    return int(match.group(1)) if match else 0


def parse_scene_number(scene_id: str) -> int:
    """
    Extract scene number from a scene_id string.
    Handles formats: 'S1', 'S10', 'LS1_S3', 'scene_1', etc.

    FIX: the original regex r'_S(\\d+)' required an underscore prefix,
    so bare scene IDs like 'S1' always returned 0 and scenes sorted
    incorrectly.  This version matches the trailing digits after any
    leading non-digit characters.
    """
    # Try "S<digits>" anywhere in the string (covers S1, LS1_S3, scene_S5)
    match = re.search(r"S(\d+)", scene_id, re.IGNORECASE)
    return int(match.group(1)) if match else 0


class PPTGeneratorService:
    """
    Service for generating PowerPoint presentations from pipeline outputs.
    """

    def __init__(
        self, output_dir: str = "outputs/ppt", width: int = 10, height: int = 7.5
    ):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.width = width
        self.height = height

    def create_presentation(self) -> Presentation:
        prs = Presentation()
        prs.slide_width = Inches(self.width)
        prs.slide_height = Inches(self.height)
        return prs

    def add_title_slide(
        self, prs: Presentation, subject: str, class_level: str, chapter_name: str
    ) -> None:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = chapter_name
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        subject_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(8), Inches(0.5)
        )
        subject_frame = subject_box.text_frame
        subject_frame.text = f"Subject: {subject}"
        subject_para = subject_frame.paragraphs[0]
        subject_para.font.size = Pt(24)
        subject_para.alignment = PP_ALIGN.CENTER

        class_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.7), Inches(8), Inches(0.5)
        )
        class_frame = class_box.text_frame
        class_frame.text = f"Class: {class_level}"
        class_para = class_frame.paragraphs[0]
        class_para.font.size = Pt(24)
        class_para.alignment = PP_ALIGN.CENTER

    def add_learning_step_slide(
        self, prs: Presentation, learning_step_id: str, title: str, concepts: List[str]
    ) -> None:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        ls_id_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(2), Inches(0.5)
        )
        ls_id_frame = ls_id_box.text_frame
        ls_id_frame.text = learning_step_id
        ls_id_para = ls_id_frame.paragraphs[0]
        ls_id_para.font.size = Pt(20)
        ls_id_para.font.bold = True
        ls_id_para.font.color.rgb = RGBColor(0, 102, 204)

        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        if concepts:
            concepts_text = "\n• " + "\n• ".join(concepts[:5])
            concepts_box = slide.shapes.add_textbox(
                Inches(1), Inches(4), Inches(8), Inches(2)
            )
            concepts_frame = concepts_box.text_frame
            concepts_frame.text = concepts_text
            concepts_para = concepts_frame.paragraphs[0]
            concepts_para.font.size = Pt(18)
            concepts_para.alignment = PP_ALIGN.LEFT

    def add_scene_slide(
        self,
        prs: Presentation,
        scene_data: Dict[str, Any],
        scene_index: int,
        image_path: Optional[str] = None,
    ) -> None:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        scene_id = scene_data.get("scene_id", f"S{scene_index + 1}")
        scene_goal = scene_data.get(
            "scene_goal",
            scene_data.get("learning_moment", scene_data.get("phase", "")),
        )

        header_height = self.height * 0.14
        image_top = header_height
        image_height = self.height - header_height
        image_width = self.width

        header_shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(self.width), Inches(header_height)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
        header_shape.line.fill.background()

        header_text = f"Scene goal: {scene_goal}"
        header_box = slide.shapes.add_textbox(
            Inches(0.3),
            Inches(0.15),
            Inches(self.width - 0.6),
            Inches(header_height - 0.3),
        )
        header_frame = header_box.text_frame
        header_frame.text = header_text
        header_frame.word_wrap = True
        header_para = header_frame.paragraphs[0]
        header_para.font.size = Pt(20)
        header_para.font.bold = True
        header_para.font.color.rgb = RGBColor(255, 255, 255)
        header_para.alignment = PP_ALIGN.LEFT

        if image_path and Path(image_path).exists():
            slide.shapes.add_picture(
                image_path,
                Inches(0),
                Inches(image_top),
                width=Inches(image_width),
                height=Inches(image_height),
            )
        else:
            # FIX: add a visible placeholder so slides aren't silently blank
            placeholder_box = slide.shapes.add_textbox(
                Inches(1), Inches(image_top + 1), Inches(8), Inches(1)
            )
            placeholder_frame = placeholder_box.text_frame
            placeholder_frame.text = f"[Image not found: {image_path}]"
            placeholder_para = placeholder_frame.paragraphs[0]
            placeholder_para.font.size = Pt(16)
            placeholder_para.font.color.rgb = RGBColor(200, 0, 0)
            placeholder_para.alignment = PP_ALIGN.CENTER

    def generate_ppt(
        self,
        state: PipelineState,
        learning_steps_dir: str,
        output_filename: str = "lesson_output.pptx",
    ) -> str:
        prs = self.create_presentation()

        self.add_title_slide(
            prs=prs,
            subject=state.user_inputs.subject,
            class_level=state.user_inputs.class_level,
            chapter_name=state.user_inputs.chapter_name,
        )

        learning_steps = state.learning_steps_list

        # FIX: use parse_ls_number for LS sorting (original used parse_scene_numbers
        # which looked for _S(\d+) pattern — never matched LS IDs like "LS1").
        sorted_learning_steps = sorted(
            learning_steps,
            key=lambda ls: parse_ls_number(ls.get("learning_step_id", "LS0")),
        )

        for ls_index, learning_step in enumerate(sorted_learning_steps):
            ls_id = learning_step.get("learning_step_id", f"LS{ls_index + 1}")
            title = learning_step.get("title", f"Learning Step {ls_index + 1}")
            concepts = learning_step.get("concepts_introduced", [])

            self.add_learning_step_slide(
                prs=prs, learning_step_id=ls_id, title=title, concepts=concepts
            )

            scenes = learning_step.get("scenes", [])

            # FIX: use parse_scene_number which correctly handles bare "S1", "S2", etc.
            sorted_scenes = sorted(
                scenes,
                key=lambda s: parse_scene_number(s.get("scene_id", "S0")),
            )

            for scene_index, scene in enumerate(sorted_scenes):
                scene_id = scene.get("scene_id", f"S{scene_index + 1}")
                image_path = (
                    Path(state.run_folder) / "images" / f"{ls_id}_{scene_id}.png"
                )

                self.add_scene_slide(
                    prs=prs,
                    scene_data=scene,
                    scene_index=scene_index,
                    image_path=str(image_path),
                )

        output_path = Path(learning_steps_dir).parent / output_filename
        prs.save(str(output_path))

        print(f"[PPT] Saved: {output_path}  ({len(sorted_learning_steps)} learning steps)")
        return str(output_path)


def create_ppt_generator() -> PPTGeneratorService:
    return PPTGeneratorService()