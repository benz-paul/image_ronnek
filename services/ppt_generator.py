"""
PPT Generator Service - Creates PowerPoint presentations from learning steps and scenes.

This service generates a final PPT file with:
- Title slide
- Learning step slides
- Scene slides for each learning step
"""

import os
import json
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from state.pipeline_state import PipelineState


class PPTGeneratorService:
    """
    Service for generating PowerPoint presentations from pipeline outputs.
    """
    
    def __init__(
        self,
        output_dir: str = "outputs/ppt",
        width: int = 10,
        height: int = 7.5
    ):
        """
        Initialize the PPT Generator Service.
        
        Args:
            output_dir: Directory to save generated PPTs
            width: Slide width in inches
            height: Slide height in inches
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        self.width = width
        self.height = height
    
    def create_presentation(self) -> Presentation:
        """
        Create a new PowerPoint presentation.
        
        Returns:
            Presentation object
        """
        prs = Presentation()
        prs.slide_width = Inches(self.width)
        prs.slide_height = Inches(self.height)
        return prs
    
    def add_title_slide(
        self,
        prs: Presentation,
        subject: str,
        class_level: str,
        chapter_name: str
    ) -> None:
        """
        Add the title slide.
        
        Slide 1:
        - Title Slide
        - Subject
        - Class
        - Chapter Name
        
        Args:
            prs: Presentation object
            subject: Subject name
            class_level: Class level
            chapter_name: Chapter name
        """
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = chapter_name
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subject
        subject_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(8), Inches(0.5)
        )
        subject_frame = subject_box.text_frame
        subject_frame.text = f"Subject: {subject}"
        subject_para = subject_frame.paragraphs[0]
        subject_para.font.size = Pt(24)
        subject_para.alignment = PP_ALIGN.CENTER
        
        # Class
        class_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.7), Inches(8), Inches(0.5)
        )
        class_frame = class_box.text_frame
        class_frame.text = f"Class: {class_level}"
        class_para = class_frame.paragraphs[0]
        class_para.font.size = Pt(24)
        class_para.alignment = PP_ALIGN.CENTER
    
    def add_learning_step_slide(
        self,
        prs: Presentation,
        learning_step_id: str,
        title: str,
        concepts: List[str]
    ) -> None:
        """
        Add a learning step title slide.
        
        Args:
            prs: Presentation object
            learning_step_id: Learning step ID (e.g., "LS1")
            title: Learning step title
            concepts: List of concepts introduced
        """
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Learning Step ID
        ls_id_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(2), Inches(0.5)
        )
        ls_id_frame = ls_id_box.text_frame
        ls_id_frame.text = learning_step_id
        ls_id_para = ls_id_frame.paragraphs[0]
        ls_id_para.font.size = Pt(20)
        ls_id_para.font.bold = True
        ls_id_para.font.color.rgb = RGBColor(0, 102, 204)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # Concepts
        if concepts:
            concepts_text = "\n• " + "\n• ".join(concepts[:5])  # Limit to 5 concepts
            concepts_box = slide.shapes.add_textbox(
                Inches(1), Inches(4), Inches(8), Inches(2)
            )
            concepts_frame = concepts_box.text_frame
            concepts_frame.text = concepts_text
            concepts_para = concepts_frame.paragraphs[0]
            concepts_para.font.size = Pt(18)
            concepts_para.alignment = PP_ALIGN.LEFT
    
    def add_scene_slide(
        self,
        prs: Presentation,
        scene_data: Dict[str, Any],
        scene_index: int
    ) -> None:
        """
        Add a scene slide with narrative and dialogue.
        
        Args:
            prs: Presentation object
            scene_data: Scene data from learning step JSON
            scene_index: Scene index number
        """
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        scene_id = scene_data.get("scene_id", f"S{scene_index + 1}")
        scene_phase = scene_data.get("scene_phase", "")
        scene_goal = scene_data.get("scene_goal", "")
        
        # Scene ID and Phase
        header_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(3), Inches(0.4)
        )
        header_frame = header_box.text_frame
        header_frame.text = f"{scene_id} - {scene_phase.title()}"
        header_para = header_frame.paragraphs[0]
        header_para.font.size = Pt(16)
        header_para.font.bold = True
        header_para.font.color.rgb = RGBColor(0, 102, 204)
        
        # Scene Goal
        goal_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.8), Inches(9), Inches(0.5)
        )
        goal_frame = goal_box.text_frame
        goal_frame.text = f"Goal: {scene_goal}"
        goal_para = goal_frame.paragraphs[0]
        goal_para.font.size = Pt(14)
        goal_para.font.italic = True
        
        # Visual Setting
        visual_setting = scene_data.get("visual_setting", {})
        environment = visual_setting.get("environment", "")
        atmosphere = visual_setting.get("atmosphere", "")
        
        setting_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.4), Inches(9), Inches(0.8)
        )
        setting_frame = setting_box.text_frame
        setting_frame.text = f"Setting: {environment}\nAtmosphere: {atmosphere}"
        setting_para = setting_frame.paragraphs[0]
        setting_para.font.size = Pt(12)
        
        # Narrative / Screenplay
        narrative = scene_data.get("narrative", {})
        screenplay = narrative.get("screenplay", "")
        
        narrative_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.3), Inches(9), Inches(2)
        )
        narrative_frame = narrative_box.text_frame
        narrative_frame.text = screenplay
        narrative_frame.word_wrap = True
        narrative_para = narrative_frame.paragraphs[0]
        narrative_para.font.size = Pt(14)
        
        # Dialogue
        dialogues = scene_data.get("dialogue", [])
        if dialogues:
            dialogue_texts = []
            for d in dialogues:
                speaker = d.get("speaker", "")
                text = d.get("text", "")
                dialogue_texts.append(f"{speaker}: {text}")
            
            dialogue_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.4), Inches(9), Inches(2)
            )
            dialogue_frame = dialogue_box.text_frame
            dialogue_frame.text = "\n".join(dialogue_texts)
            dialogue_frame.word_wrap = True
            
            for para in dialogue_frame.paragraphs:
                para.font.size = Pt(12)
                if para.text.startswith("Aarav") or para.text.startswith("Meera"):
                    para.font.bold = True
    
    def generate_ppt(
        self,
        state: PipelineState,
        learning_steps_dir: str,
        output_filename: str = "lesson_output.pptx"
    ) -> str:
        """
        Generate the complete PowerPoint presentation.
        
        Slide Structure:
        - Slide 1: Title Slide (Subject, Class, Chapter Name)
        - Slide 2: Learning Step 1 Title
        - Next Slides: All scenes from Learning Step 1
        - Next Slide: Learning Step 2 Title
        - Next Slides: All scenes from Learning Step 2
        - Continue until all learning steps included
        
        Args:
            state: Current pipeline state
            learning_steps_dir: Directory containing learning step JSON files
            output_filename: Output filename
            
        Returns:
            Path to generated PPT file
        """
        prs = self.create_presentation()
        
        # Add title slide
        self.add_title_slide(
            prs=prs,
            subject=state.user_inputs.subject,
            class_level=state.user_inputs.class_level,
            chapter_name=state.user_inputs.chapter_name
        )
        
        # Process each learning step
        learning_steps = state.learning_steps_list
        
        for ls_index, learning_step in enumerate(learning_steps):
            # Get learning step details
            ls_id = learning_step.get("learning_step_id", f"LS{ls_index + 1}")
            title = learning_step.get("title", f"Learning Step {ls_index + 1}")
            concepts = learning_step.get("concepts_introduced", [])
            
            # Add learning step slide
            self.add_learning_step_slide(
                prs=prs,
                learning_step_id=ls_id,
                title=title,
                concepts=concepts
            )
            
            # Add scene slides
            scenes = learning_step.get("scenes", [])
            for scene_index, scene in enumerate(scenes):
                self.add_scene_slide(
                    prs=prs,
                    scene_data=scene,
                    scene_index=scene_index
                )
        
        # Save the presentation
        output_path = Path(learning_steps_dir).parent / output_filename
        prs.save(str(output_path))
        
        return str(output_path)


def create_ppt_generator() -> PPTGeneratorService:
    """Factory function to create PPTGeneratorService."""
    return PPTGeneratorService()
